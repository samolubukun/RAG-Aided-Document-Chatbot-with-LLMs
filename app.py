import streamlit as st
import os
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain_community.llms import Cohere
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import PromptTemplate
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Set API keys
os.environ['HuggingFaceHub_API_Token'] = os.getenv('HUGGINGFACEHUB_API_TOKEN')
os.environ['GOOGLE_API_KEY'] = os.getenv('GOOGLE_API_KEY')
os.environ['cohere_api_key'] = os.getenv('COHERE_API_KEY')

# Initialize session state
if 'vectorstore' not in st.session_state:
    st.session_state.vectorstore = None

# Function to extract text from PDF
def extract_text_from_pdf(pdf_file):
    pdf_text = ""
    pdf_reader = PdfReader(pdf_file)
    for page in pdf_reader.pages:
        pdf_text += page.extract_text()
    return pdf_text

# Function to extract text from PPTX
def extract_text_from_pptx(pptx_file):
    ppt_text = ""
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                ppt_text += shape.text + '\n'
    return ppt_text

# Function to extract text from DOCX
def extract_text_from_docx(docx_file):
    doc_text = ""
    doc = Document(docx_file)
    for paragraph in doc.paragraphs:
        doc_text += paragraph.text + '\n'
    return doc_text

# Function to process uploaded documents
def process_documents(uploaded_files):
    all_text = ""
    for file in uploaded_files:
        if file.name.endswith('.pdf'):
            all_text += extract_text_from_pdf(file)
        elif file.name.endswith('.pptx'):
            all_text += extract_text_from_pptx(file)
        elif file.name.endswith('.docx'):
            all_text += extract_text_from_docx(file)

    # Split text into chunks
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
        separators=['\n', '\n\n', ' ', '']
    )
    chunks = text_splitter.split_text(text=all_text)

    # Create embeddings and vectorstore
    embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
    vectorstore = FAISS.from_texts(chunks, embedding=embeddings)
    return vectorstore

# Function to format documents
def format_docs(docs):
    return "\n\n".join(doc.page_content for doc in docs)

# Function to generate answers
def generate_answer(question, vectorstore):
    retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 6})
    
    prompt_template = """Answer the question as precisely as possible using the provided context. If the answer is
                         not contained in the context, say "answer not available in context" \n\n
                         Context: \n {context}?\n
                         Question: \n {question} \n
                         Answer:"""
    
    prompt = PromptTemplate.from_template(template=prompt_template)
    
    cohere_llm = Cohere(model="command", temperature=0.1, cohere_api_key=os.getenv('cohere_api_key'))
    
    rag_chain = (
        {"context": retriever | format_docs, "question": RunnablePassthrough()}
        | prompt
        | cohere_llm
        | StrOutputParser()
    )
    
    return rag_chain.invoke(question)

# Function to generate MCQs
def generate_mcqs(num_questions, vectorstore):
    retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 6})

    mcq_prompt_template = """Generate {num_questions} multiple-choice questions based on the provided context. 
    Each question should have 4 options, with one correct answer clearly indicated.

    Context: 
    {context}

    Questions:
    """

    prompt = PromptTemplate.from_template(template=mcq_prompt_template)

    cohere_llm = Cohere(model="command", temperature=0.7, cohere_api_key=os.getenv('cohere_api_key'))

    rag_chain = (
        {"context": retriever | format_docs, "num_questions": RunnablePassthrough()}
        | prompt
        | cohere_llm
        | StrOutputParser()
    )

    return rag_chain.invoke(num_questions)

# Streamlit UI
st.title("RAG-Aided Document Chatbot with LLMs")

# File upload section
uploaded_files = st.file_uploader("Upload your documents (PDF, PPTX, DOCX)", 
                                  type=['pdf', 'pptx', 'docx'],
                                  accept_multiple_files=True)

# Process documents button with unique key
if uploaded_files and st.button("Process Documents", key="process_docs"):
    with st.spinner("Processing documents..."):
        st.session_state.vectorstore = process_documents(uploaded_files)
    st.success("Documents processed successfully!")

# Question input for answering questions
question = st.text_input("Ask a question about your documents:")

# Generate answer based on question
if question and st.session_state.vectorstore is not None:
    with st.spinner("Generating answer..."):
        answer = generate_answer(question, st.session_state.vectorstore)
    st.write("Answer:", answer)
elif question:
    st.warning("Please upload and process documents first.")

# MCQ generation section
num_questions = st.number_input("Enter the number of MCQs to generate:", min_value=1, max_value=20, step=1)

if num_questions and st.session_state.vectorstore is not None and st.button("Generate MCQs", key="generate_mcqs"):
    with st.spinner("Generating MCQs..."):
        mcqs = generate_mcqs(num_questions, st.session_state.vectorstore)
    st.write("Generated MCQs:")
    st.write(mcqs)
elif num_questions:
    st.warning("Please upload and process documents first.")
